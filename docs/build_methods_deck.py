# Build a methods-focused slide deck for the MN proxy-prediction project.
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- palette (Ocean / scientific) ----
NAVY   = RGBColor(0x1B, 0x2A, 0x4A)   # dark bg
BLUE   = RGBColor(0x06, 0x5A, 0x82)   # primary
TEAL   = RGBColor(0x1C, 0x72, 0x93)   # secondary
TEAL2  = RGBColor(0x2E, 0x8B, 0xA6)
AMBER  = RGBColor(0xE3, 0x9A, 0x2E)   # accent / problem
GREEN  = RGBColor(0x3F, 0x8E, 0x6B)   # works
RED    = RGBColor(0xB8, 0x4A, 0x3E)   # doesn't
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
INK    = RGBColor(0x1E, 0x2A, 0x38)   # body text
MUTE   = RGBColor(0x5A, 0x6B, 0x7B)
LIGHT  = RGBColor(0xEE, 0xF3, 0xF6)   # card bg
CARDLN = RGBColor(0xD3, 0xDE, 0xE6)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = 13.333, 7.5
BLANK = prs.slide_layouts[6]
HEAD = "Georgia"
BODY = "Calibri"

def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.line.fill.background(); r.fill.solid(); r.fill.fore_color.rgb = bg
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return s

def box(s, x, y, w, h, fill=None, line=None, rounded=False, lw=1.0):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shp.shadow.inherit = False
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(lw)
    return shp

def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=2):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2); tf.margin_top = tf.margin_bottom = Pt(2)
    if isinstance(runs, str): runs = [(runs, {})]
    if runs and isinstance(runs[0], tuple): runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space); p.space_before = Pt(0)
        for t, o in para:
            r = p.add_run(); r.text = t
            f = r.font; f.name = o.get("f", BODY); f.size = Pt(o.get("s", 14))
            f.bold = o.get("b", False); f.italic = o.get("i", False)
            f.color.rgb = o.get("c", INK)
        if "bullet" in (para[0][1] if para else {}):
            pass
    return tb

def bullets(s, x, y, w, h, items, sz=14, color=INK, gap=5, mk="—", mkc=None):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(2); tf.margin_top = Pt(2)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.space_before = Pt(0)
        if isinstance(it, tuple): head, rest = it
        else: head, rest = None, it
        rm = p.add_run(); rm.text = mk + "  "
        rm.font.name = BODY; rm.font.size = Pt(sz); rm.font.bold = True
        rm.font.color.rgb = mkc or TEAL
        if head:
            rh = p.add_run(); rh.text = head + ": "
            rh.font.name = BODY; rh.font.size = Pt(sz); rh.font.bold = True; rh.font.color.rgb = color
        rr = p.add_run(); rr.text = rest
        rr.font.name = BODY; rr.font.size = Pt(sz); rr.font.color.rgb = color
    return tb

def title_bar(s, kicker, ttl, num):
    text(s, 0.6, 0.42, 11, 0.3, [[(kicker, {"f":BODY,"s":12.5,"b":True,"c":TEAL})]])
    text(s, 0.6, 0.72, 11.6, 0.9, [[(ttl, {"f":HEAD,"s":29,"b":True,"c":NAVY})]])
    text(s, 12.5, 0.5, 0.6, 0.4, [[(num, {"f":BODY,"s":12,"b":True,"c":CARDLN})]], align=PP_ALIGN.RIGHT)

# ============================================================ 1. TITLE
s = slide(NAVY)
box(s, 0, 0, 0.28, SH, fill=AMBER)
text(s, 0.9, 1.9, 11.6, 0.4, [[("MICRONUTRIENT DEFICIENCY PREDICTION FROM PROXY INDICATORS", {"f":BODY,"s":13,"b":True,"c":TEAL2})]])
text(s, 0.9, 2.35, 11.6, 2.0, [
    [("Methods We've Tried — and the", {"f":HEAD,"s":40,"b":True,"c":WHITE})],
    [("Transportability Problem", {"f":HEAD,"s":40,"b":True,"c":AMBER})]])
text(s, 0.9, 4.35, 11.6, 0.9, [[("Feature engineering, estimation strategies, and why cross-country prediction is hard", {"f":BODY,"s":16,"c":RGBColor(0xC6,0xD4,0xDE)})]])
box(s, 0.95, 5.55, 3.4, 0.02, fill=TEAL)
text(s, 0.9, 5.75, 11.6, 0.5, [[("Gambia · Ghana · Sierra Leone · Malawi   |   methods overview", {"f":BODY,"s":12.5,"c":RGBColor(0x9F,0xB2,0xC2)})]])

# ============================================================ 2. THE CORE CHALLENGE
s = slide()
title_bar(s, "THE SETUP", "One challenge shapes every method choice", "02")
# left narrative
bullets(s, 0.6, 1.85, 6.4, 4.6, [
    ("Goal", "predict micronutrient-deficiency prevalence for every Admin-2 district, including unsurveyed ones, from widely-available proxies."),
    ("Predictors", "1,000+ proxy variables — DHS/MICS indicators, satellite (GEE), climate, malaria, soil, food prices."),
    ("Biomarkers", "iron, vitamin A, folate, B12, zinc from 4 national surveys (2013–2018)."),
    ("The catch", "proxies live at Admin-2 level, so the EFFECTIVE sample size is the number of AREAS, not individuals."),
], sz=15, gap=11)
# right: the n callout
box(s, 7.4, 1.95, 5.3, 4.25, fill=LIGHT, line=CARDLN, rounded=True, lw=1)
text(s, 7.7, 2.25, 4.8, 0.4, [[("EFFECTIVE SAMPLE SIZE", {"f":BODY,"s":12.5,"b":True,"c":TEAL})]])
text(s, 7.7, 2.7, 4.8, 1.0, [[("14–87", {"f":HEAD,"s":58,"b":True,"c":BLUE})]])
text(s, 7.7, 3.75, 4.8, 0.4, [[("Admin-2 areas per country", {"f":BODY,"s":13.5,"c":MUTE})]])
box(s, 7.7, 4.35, 4.7, 0.012, fill=CARDLN)
text(s, 7.7, 4.55, 4.9, 1.5, [
    [("p ≫ n.", {"f":BODY,"s":15,"b":True,"c":AMBER}), (" 1,000+ predictors, ~dozens of independent observations — the root cause of weak, unstable, hard-to-transport prediction.", {"f":BODY,"s":14,"c":INK})]])

# ============================================================ 3. FEATURE ENGINEERING
s = slide()
title_bar(s, "WHAT WE'VE TRIED · 1 of 2", "Feature engineering", "03")
cards = [
    ("Predictor scope", TEAL, ["Proxy-only (exclude survey gw_ vars)", "Domains matched to outcome biology"]),
    ("Dimensionality", BLUE, ["Within-domain PCA", "Greedy decorrelation (top-k |r|)", "5-stage prescreening"]),
    ("Transforms", TEAL2, ["Rank-normalization (not z-score)", "Median imputation"]),
    ("Variable selection", BLUE, ["DAG-based selection", "Forward selection", "Invariance / combined filters"]),
    ("Constructs", TEAL, ["Outcome-aware interactions (fe_*)", "Transportable constructs", "EconomicClusters wealth rank (planned)"]),
    ("External features", TEAL2, ["SoilGrids (+0.045 LOCO r)", "GEE · CHIRPS · MAP · IHME · WFP"]),
]
cx, cy, cw, ch, gx, gy = 0.6, 1.8, 3.95, 1.62, 0.27, 0.26
for i, (h, col, items) in enumerate(cards):
    x = cx + (i % 3) * (cw + gx); y = cy + (i // 3) * (ch + gy)
    box(s, x, y, cw, ch, fill=LIGHT, line=CARDLN, rounded=True)
    box(s, x, y, 0.12, ch, fill=col)
    text(s, x + 0.28, y + 0.13, cw - 0.4, 0.35, [[(h, {"f":BODY,"s":14.5,"b":True,"c":NAVY})]])
    bullets(s, x + 0.28, y + 0.55, cw - 0.45, ch - 0.6, items, sz=11, gap=2.5, mk="·", mkc=col)
text(s, 0.6, 6.7, 12, 0.4, [[("Across all of it the binding constraint persists — feature work sharpens within-country signal but cannot manufacture cross-country transport.", {"f":BODY,"s":12.5,"i":True,"c":MUTE})]])

# ============================================================ 4. ESTIMATION
s = slide()
title_bar(s, "WHAT WE'VE TRIED · 2 of 2", "Estimation strategies", "04")
groups = [
    ("Individual-level SuperLearner", BLUE, "6-learner stack, aggregated to Admin-2  ·  now a SENSITIVITY analysis"),
    ("Area-level SuperLearner", GREEN, "fit directly on Admin-2 aggregates (MSE & NLL loss)  ·  now PRIMARY"),
    ("Bayesian small-area (SAE)", TEAL, "Fay–Herriot  ·  BYM2 spatial (INLA, PC priors)"),
    ("Frequentist regressions", TEAL2, "OLS/GLM · elastic-net · group-lasso · GAM · mixed · quantile · beta · quasi-binomial"),
    ("Spatial models", TEAL, "coordinates · spatial GAM · spatial+SoilGrids · SPDE-INLA · Gaussian process"),
    ("Ensembles & hybrids", BLUE, "NNLS stacker · two-stage (SL→spatial GAM) · SL→BYM2 (calibrated map intervals)"),
    ("Calibration & uncertainty", TEAL2, "conformal prediction intervals · Platt recalibration"),
    ("Resolution variants", TEAL, "Admin-2 (primary) · Admin-3 (Sierra Leone chiefdoms) · survey-cluster GPS buffers"),
]
gy0 = 1.8; gh = 0.585; gp = 0.04
for i, (h, col, d) in enumerate(groups):
    y = gy0 + i * (gh + gp)
    box(s, 0.6, y, 12.13, gh, fill=LIGHT if i % 2 == 0 else WHITE, line=CARDLN, rounded=False)
    box(s, 0.6, y, 0.1, gh, fill=col)
    text(s, 0.85, y, 3.95, gh, [[(h, {"f":BODY,"s":13.5,"b":True,"c":NAVY})]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 4.9, y, 7.7, gh, [[(d, {"f":BODY,"s":12.5,"c":INK})]], anchor=MSO_ANCHOR.MIDDLE)

# ============================================================ 5. PRIMARY VS SENSITIVITY
s = slide()
title_bar(s, "WHAT CHANGED", "From one model to a primary + sensitivity split", "05")
# primary card
box(s, 0.6, 1.95, 5.95, 4.2, fill=LIGHT, line=GREEN, rounded=True, lw=1.5)
box(s, 0.6, 1.95, 5.95, 0.7, fill=GREEN, rounded=True)
text(s, 0.85, 1.97, 5.5, 0.66, [[("PRIMARY  —  Admin-2 area-level SAE", {"f":BODY,"s":15.5,"b":True,"c":WHITE})]], anchor=MSO_ANCHOR.MIDDLE)
bullets(s, 0.9, 2.9, 5.4, 3.1, [
    "Models fit directly at Admin-2 on survey-weighted prevalence",
    "Area-level SuperLearner, Fay–Herriot, SL→BYM2",
    "Strong within-country ranking; full district coverage",
    "Native, spatially-smoothed credible intervals",
    "Headline results + dashboard default",
], sz=14, gap=9, mkc=GREEN)
# sensitivity card
box(s, 6.85, 1.95, 5.88, 4.2, fill=LIGHT, line=MUTE, rounded=True, lw=1.5)
box(s, 6.85, 1.95, 5.88, 0.7, fill=MUTE, rounded=True)
text(s, 7.1, 1.97, 5.5, 0.66, [[("SENSITIVITY  —  individual-level SL", {"f":BODY,"s":15.5,"b":True,"c":WHITE})]], anchor=MSO_ANCHOR.MIDDLE)
bullets(s, 7.15, 2.9, 5.35, 3.1, [
    "Person-level SuperLearner aggregated to Admin-2",
    "Aggregating noisy individual predictions adds bias",
    "Retained for comparison, relocated to R/sensitivity/",
    "Not the headline estimator",
], sz=14, gap=9, mkc=MUTE)
text(s, 0.6, 6.45, 12, 0.5, [[("Why: ", {"f":BODY,"s":13,"b":True,"c":NAVY}), ("with so few areas, SAE methods match or beat the flexible ensemble — and the BMGF call endorsed leading with SAE and feeding SL predictions into it.", {"f":BODY,"s":13,"i":True,"c":MUTE})]])

# ============================================================ 6. THE CORE TRANSPORTABILITY ISSUE
s = slide(NAVY)
box(s, 0, 0, SW, 1.45, fill=BLUE)
text(s, 0.6, 0.32, 11, 0.3, [[("THE CORE ISSUE", {"f":BODY,"s":12.5,"b":True,"c":AMBER})]])
text(s, 0.6, 0.62, 12, 0.7, [[("Within country it works. Across countries it doesn't.", {"f":HEAD,"s":27,"b":True,"c":WHITE})]])
stats = [
    ("0.02", "median LOCO r\nacross 16 folds", AMBER),
    ("≈0.1–0.2", "per-method LOCO r\n(near chance)", WHITE),
    ("SAE ≥ SL", "on most leave-one-\ncountry-out splits", WHITE),
    ("14%", "of error removed by\nfixing the level offset", AMBER),
]
for i, (big, lab, col) in enumerate(stats):
    x = 0.6 + i * 3.12
    box(s, x, 1.95, 2.9, 1.85, fill=RGBColor(0x25,0x37,0x5C), line=None, rounded=True)
    text(s, x+0.1, 2.15, 2.7, 0.9, [[(big, {"f":HEAD,"s":33,"b":True,"c":col})]], align=PP_ALIGN.CENTER)
    text(s, x+0.1, 3.05, 2.7, 0.7, [[(lab.replace("\n"," "), {"f":BODY,"s":12.5,"c":RGBColor(0xC6,0xD4,0xDE)})]], align=PP_ALIGN.CENTER)
bullets(s, 0.7, 4.25, 12, 2.7, [
    ("Within-country, area-level", "ranks Admin-2 districts well — but the headline r > 0.9 is in-sample; genuine out-of-fold r ≈ 0.25–0.53 (a flagged bug)."),
    ("Cross-country (LOCO)", "the SuperLearner is among the WEAKEST methods; flexibility doesn't pay when n = areas."),
    ("The failure is mostly LEVEL, not ranking", "models often rank a held-out country's districts ~ok but miss the absolute prevalence by tens of points."),
], sz=14.5, color=RGBColor(0xDD,0xE6,0xEE), gap=10, mkc=AMBER)

# ============================================================ 7. FOUR FAILURE MODES
s = slide()
title_bar(s, "DIAGNOSIS", "Transport fails for four distinct reasons", "07")
fm = [
    ("1 · Gambia iron", AMBER, "Real biology", "Raw ferritin genuinely ~6× lower; 66.6% deficiency corroborated by external BRINDA analysis. A true level gap — NOT fixable by modelling."),
    ("2 · Folate / B12", GREEN, "Data artifact (fixed)", "Wrong units + cutoff made Malawi folate read ~0%. Corrected to ~21% at WHO <10 nmol/L, validated vs published survey."),
    ("3 · Malawi / SL iron", RED, "Pattern non-transfer", "Even district RANKING fails (r ≈ 0) — a genuine West- vs East-Africa difference models trained elsewhere can't recover."),
    ("4 · Vitamin A (both)", TEAL, "Weak signal", "RBP is well-harmonized, but proxies simply don't predict the VAD spatial pattern (r ≈ 0–0.25); women's VAD is near-constant."),
]
cw, ch, gx, gy = 5.95, 2.18, 0.23, 0.24
for i, (h, col, tag, d) in enumerate(fm):
    x = 0.6 + (i % 2) * (cw + gx); y = 1.85 + (i // 2) * (ch + gy)
    box(s, x, y, cw, ch, fill=LIGHT, line=CARDLN, rounded=True)
    box(s, x, y, cw, 0.55, fill=col, rounded=True)
    text(s, x+0.25, y, cw-1.7, 0.55, [[(h, {"f":BODY,"s":14.5,"b":True,"c":WHITE})]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, x+cw-1.75, y, 1.6, 0.55, [[(tag, {"f":BODY,"s":11.5,"b":True,"c":WHITE})]], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
    text(s, x+0.25, y+0.7, cw-0.5, ch-0.85, [[(d, {"f":BODY,"s":12.5,"c":INK})]])

# ============================================================ 8. BIOMARKER HARMONIZATION
s = slide()
title_bar(s, "THE LEADING SUSPECT", "The same outcome is built four different ways", "08")
# table-like comparison
hdr = ["", "Iron continuous", "Adjustment", "Scale"]
rows = [
    ["Gambia", "gw_LogFerAdj", "BRINDA", "log"],
    ["Ghana", "gw_cFerrAdjThurn", "Thurnham", "linear"],
    ["Sierra Leone", "gw_cFerrAdj", "unspecified", "linear"],
    ["Malawi", "sf_reg", "survey-native", "linear"],
]
tx, ty, colw = 0.6, 1.95, [2.2, 3.1, 2.6, 1.4]
# header
cxp = tx
for j, htxt in enumerate(hdr):
    box(s, cxp, ty, colw[j], 0.5, fill=NAVY)
    text(s, cxp+0.12, ty, colw[j]-0.2, 0.5, [[(htxt, {"f":BODY,"s":12.5,"b":True,"c":WHITE})]], anchor=MSO_ANCHOR.MIDDLE)
    cxp += colw[j]
for r_i, row in enumerate(rows):
    cxp = tx; yy = ty + 0.5 + r_i*0.5
    for j, cell in enumerate(row):
        box(s, cxp, yy, colw[j], 0.5, fill=(LIGHT if r_i%2==0 else WHITE), line=CARDLN)
        bold = (j==0)
        col = AMBER if (j==2 and r_i<2) else INK
        text(s, cxp+0.12, yy, colw[j]-0.2, 0.5, [[(cell, {"f":(BODY),"s":12,"b":bold,"c":col})]], anchor=MSO_ANCHOR.MIDDLE)
        cxp += colw[j]
# right insight panel
box(s, 9.7, 1.95, 3.03, 2.5, fill=NAVY, rounded=True)
text(s, 9.95, 2.2, 2.55, 0.4, [[("WHY IT BREAKS TRANSPORT", {"f":BODY,"s":11.5,"b":True,"c":AMBER})]])
text(s, 9.95, 2.7, 2.6, 1.7, [[("A model trained on three adjustment methods can't predict a fourth country's level. But adjustment harmonization is NOT the fix — the gap is in the RAW biomarker.", {"f":BODY,"s":12,"c":WHITE})]])
bullets(s, 0.6, 4.75, 12, 2.0, [
    ("Folate / B12", "different units, cutoffs and assays across surveys — partly genuine data bugs (now fixed for folate)."),
    ("The honest takeaway", "report transport on RANK scales alongside absolute, and state plainly that absolute prevalence does not transport across countries."),
], sz=14, gap=9, mkc=TEAL)

# ============================================================ 9. WHAT WORKS / DOESN'T
s = slide()
title_bar(s, "HONEST SUMMARY", "What works, what doesn't", "09")
box(s, 0.6, 1.95, 5.95, 4.4, fill=RGBColor(0xEC,0xF4,0xEF), line=GREEN, rounded=True, lw=1.5)
text(s, 0.9, 2.2, 5.4, 0.5, [[("WORKS", {"f":BODY,"s":16,"b":True,"c":GREEN})]])
bullets(s, 0.9, 2.85, 5.4, 3.4, [
    "Within-country Admin-2 maps (area-level SAE)",
    "Full-coverage prediction for unsurveyed districts",
    "Calibrated credible intervals (SL→BYM2)",
    "National prevalence estimates",
    "Identifying within-country hotspots vs homogeneity",
], sz=14.5, gap=11, mkc=GREEN)
box(s, 6.85, 1.95, 5.88, 4.4, fill=RGBColor(0xF6,0xEC,0xEA), line=RED, rounded=True, lw=1.5)
text(s, 7.15, 2.2, 5.4, 0.5, [[("DOESN'T (YET)", {"f":BODY,"s":16,"b":True,"c":RED})]])
bullets(s, 7.15, 2.85, 5.35, 3.4, [
    "Cross-country transport of absolute prevalence",
    "The ensemble beating standard SAE",
    "Vitamin A spatial prediction (weak signal)",
    "Comparing biomarker levels across surveys/years",
], sz=14.5, gap=11, mkc=RED, color=INK)
text(s, 0.6, 6.55, 12, 0.5, [[("Net: a rigorous, honestly-bounded subnational tool — not a cross-country deficiency predictor.", {"f":BODY,"s":13.5,"i":True,"c":MUTE})]])

# ============================================================ 10. CONCLUSIONS
s = slide(NAVY)
box(s, 0, 0, 0.28, SH, fill=AMBER)
text(s, 0.9, 0.7, 11, 0.4, [[("WHERE THIS LEAVES US", {"f":BODY,"s":13,"b":True,"c":TEAL2})]])
text(s, 0.9, 1.1, 11.6, 0.7, [[("Conclusions & directions", {"f":HEAD,"s":30,"b":True,"c":WHITE})]])
bullets(s, 0.9, 2.15, 11.7, 3.0, [
    ("Lead with area-level SAE", "the defensible primary deliverable; individual-level SL is a sensitivity analysis."),
    ("Be honest about transport", "cross-country absolute prevalence is not deployable; report rank + intervals, not point levels."),
    ("Harmonize biomarkers", "folate units/cutoff fixed; next reconcile B12 assays and flag Gambia's real iron burden."),
    ("Combine, don't replace", "feed SuperLearner predictions into SAE (SL→BYM2) for calibrated within-country maps."),
], sz=15.5, color=RGBColor(0xDD,0xE6,0xEE), gap=13, mkc=AMBER)
box(s, 0.9, 5.5, 11.5, 1.25, fill=RGBColor(0x25,0x37,0x5C), rounded=True)
text(s, 1.2, 5.62, 11, 1.0, [[("Next external work: ", {"f":BODY,"s":14,"b":True,"c":AMBER}), ("GBD / VMNIS covariate comparison · ML-assisted SAE methods · cross-survey biomarker calibration · EconomicClusters wealth proxy.", {"f":BODY,"s":14,"c":WHITE})]], anchor=MSO_ANCHOR.MIDDLE)

prs.save(r"C:\Users\andre\OneDrive\Documents\mn-prediction\docs\MN_methods_overview.pptx")
print("saved docs/MN_methods_overview.pptx ; slides:", len(prs.slides._sldIdLst))
